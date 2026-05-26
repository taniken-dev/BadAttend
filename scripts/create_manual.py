"""
BadAttend システムマニュアル .pptx 生成スクリプト
フォント: MSPゴシック
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ===================== カラーパレット =====================
C_NAV      = RGBColor(0x1A, 0x2E, 0x4A)   # ネイビー（ヘッダー）
C_ACCENT   = RGBColor(0x25, 0x7A, 0xBF)   # アクセントブルー
C_BG       = RGBColor(0xF4, 0xF7, 0xFB)   # 背景（淡いグレー）
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x1E, 0x2D, 0x3D)   # 本文テキスト
C_MUTED    = RGBColor(0x6C, 0x7A, 0x8A)   # サブテキスト
C_LINE     = RGBColor(0xD9, 0xE4, 0xF0)   # 区切り線

# ロールカラー
C_MEMBER   = RGBColor(0x1E, 0x8B, 0x4C)   # 部員: グリーン
C_MANAGER  = RGBColor(0xD4, 0x7F, 0x00)   # マネージャー: ゴールド
C_ADMIN    = RGBColor(0x7B, 0x2D, 0x9E)   # 幹部: パープル
C_COACH    = RGBColor(0x4A, 0x6A, 0x8A)   # 顧問: スレート

# スクショプレースホルダー色
C_SS_BG    = RGBColor(0xE8, 0xEE, 0xF5)
C_SS_BORDER= RGBColor(0xA0, 0xB5, 0xC8)

FONT_NAME  = "MS PGothic"
SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # 完全ブランク


# ===================== ユーティリティ =====================

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=14, bold=False, color=C_DARK,
                 align=PP_ALIGN.LEFT, wrap=True,
                 line_spacing=None, space_before=None, space_after=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    if line_spacing:
        from pptx.util import Pt as _Pt
        para.line_spacing = line_spacing
    if space_before is not None:
        para.space_before = Pt(space_before)
    if space_after is not None:
        para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_bullet_para(tf, text, level=0, font_size=13, color=C_DARK, bold=False,
                    space_before=4, space_after=2, bullet_char="●"):
    from pptx.util import Pt as _Pt
    para = tf.add_paragraph()
    para.level = level
    if space_before:
        para.space_before = Pt(space_before)
    if space_after:
        para.space_after = Pt(space_after)
    run = para.add_run()
    prefix = ("　" * level) + (bullet_char + " " if level == 0 else "　 ・")
    run.text = prefix + text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return para


def add_screenshot_placeholder(slide, x, y, w, h, label="【 スクリーンショット貼付欄 】"):
    """スクリーンショット用プレースホルダー（枠＋テキスト）"""
    rect = add_rect(slide, x, y, w, h, fill=C_SS_BG, line=C_SS_BORDER, line_w=Pt(1.5))
    # ダッシュ枠をエミュレート（実線）
    add_text_box(slide, label, x, y + h/2 - Inches(0.2), w, Inches(0.4),
                 font_size=12, color=C_SS_BORDER, align=PP_ALIGN.CENTER)
    return rect


def add_role_badge(slide, role_label, color, x, y, w=Inches(1.8), h=Inches(0.38)):
    rect = add_rect(slide, x, y, w, h, fill=color)
    add_text_box(slide, role_label, x, y, w, h,
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    return rect


# ===================== 共通レイアウト部品 =====================

def draw_base(slide, bg_color=C_BG):
    """背景色"""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=bg_color)


def draw_header(slide, title, subtitle=None, role_color=C_ACCENT):
    """上部ヘッダーバー（高さ固定 1.1in）"""
    H = Inches(1.1)
    add_rect(slide, 0, 0, SLIDE_W, H, fill=C_NAV)
    # アクセントライン
    add_rect(slide, 0, H - Pt(4), SLIDE_W, Pt(4), fill=role_color)

    add_text_box(slide, title,
                 Inches(0.4), Inches(0.12), Inches(11), Inches(0.58),
                 font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.65), Inches(9), Inches(0.38),
                     font_size=13, color=RGBColor(0xA0,0xC0,0xE0), align=PP_ALIGN.LEFT)


def draw_footer(slide, page_num, total=None, note=None):
    """下部フッター"""
    Y = SLIDE_H - Inches(0.38)
    add_rect(slide, 0, Y, SLIDE_W, Inches(0.38), fill=C_NAV)
    label = f"千葉工業大学バドミントン部 │ BadAttend システムマニュアル"
    add_text_box(slide, label, Inches(0.3), Y + Pt(4), Inches(9), Inches(0.3),
                 font_size=9, color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.LEFT)
    pn = f"P. {page_num}" + (f" / {total}" if total else "")
    add_text_box(slide, pn, Inches(11.5), Y + Pt(4), Inches(1.5), Inches(0.3),
                 font_size=9, color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.RIGHT)
    if note:
        add_text_box(slide, note, Inches(0.3), Y - Inches(0.28), Inches(12.5), Inches(0.28),
                     font_size=9, color=C_MUTED, align=PP_ALIGN.LEFT)


def content_y():
    return Inches(1.25)


def content_h():
    return SLIDE_H - Inches(1.25) - Inches(0.45)


# ===================== スライド定義 =====================

# --- スライド 1: タイトル ---
def slide_title():
    sl = prs.slides.add_slide(BLANK)
    # グラデーション風の背景
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=C_NAV)
    add_rect(sl, 0, 0, SLIDE_W, Inches(4.5), fill=RGBColor(0x12,0x20,0x35))

    # 装飾ライン
    add_rect(sl, Inches(0.5), Inches(2.0), Inches(0.06), Inches(2.4), fill=C_ACCENT)
    add_rect(sl, 0, Inches(5.5), SLIDE_W, Pt(3), fill=C_ACCENT)

    add_text_box(sl, "BadAttend", Inches(0.8), Inches(1.3), Inches(12), Inches(1.0),
                 font_size=54, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text_box(sl, "システムマニュアル", Inches(0.8), Inches(2.3), Inches(12), Inches(0.65),
                 font_size=30, bold=False, color=RGBColor(0x80,0xB8,0xE8), align=PP_ALIGN.LEFT)
    add_text_box(sl, "千葉工業大学 バドミントン部 出欠管理システム",
                 Inches(0.8), Inches(3.1), Inches(12), Inches(0.5),
                 font_size=16, color=RGBColor(0x70,0xA0,0xC8), align=PP_ALIGN.LEFT)

    # バージョン・日付エリア
    add_text_box(sl, "Ver. 1.0　│　2026年5月作成",
                 Inches(0.8), Inches(5.8), Inches(8), Inches(0.4),
                 font_size=12, color=RGBColor(0x60,0x80,0xA0), align=PP_ALIGN.LEFT)

    # 凡例バッジ
    roles = [
        ("部　員", C_MEMBER),
        ("マネージャー", C_MANAGER),
        ("幹　部", C_ADMIN),
        ("顧　問", C_COACH),
    ]
    for i, (r, c) in enumerate(roles):
        add_role_badge(sl, r, c,
                       Inches(0.8) + Inches(2.1)*i, Inches(6.5),
                       Inches(1.9), Inches(0.52))

    add_text_box(sl, "本マニュアルはロール別の色分けバッジで対象者を示します",
                 Inches(0.8), Inches(7.1), Inches(10), Inches(0.3),
                 font_size=10, color=RGBColor(0x50,0x70,0x90))
    return sl


# --- スライド 2: 目次 ---
def slide_toc():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "目　次", "Contents")
    draw_footer(sl, 2)

    Y = content_y()
    col_w = Inches(5.8)
    col_gap = Inches(0.5)
    col1_x = Inches(0.5)
    col2_x = col1_x + col_w + col_gap

    items_col1 = [
        ("1", "システム概要・目的",        "P.3"),
        ("2", "ログイン方法",              "P.4"),
        ("3", "ロール体系",                "P.5"),
        ("4", "ダッシュボード",            "P.6"),
        ("5", "カレンダー・出欠登録",       "P.7"),
        ("6", "出欠連絡のルール",          "P.8"),
        ("7", "ポイントシステム",          "P.9"),
        ("8", "ランキング",               "P.10"),
    ]
    items_col2 = [
        ("9",  "マネージャー機能",         "P.11"),
        ("10", "幹部：メンバー管理",        "P.12"),
        ("11", "幹部：実績確定",           "P.13"),
        ("12", "幹部：警告フラグ",         "P.14"),
        ("13", "選考スコアの仕組み",        "P.15"),
        ("14", "体調不良ロック",           "P.16"),
        ("15", "よくある質問（FAQ）",      "P.17"),
        ("16", "お問い合わせ",            "P.18"),
    ]

    def draw_toc_col(items, x):
        for i, (num, label, pg) in enumerate(items):
            row_y = Y + Inches(0.55) * i
            row_h = Inches(0.48)
            bg = C_WHITE if i % 2 == 0 else RGBColor(0xEC,0xF3,0xFA)
            add_rect(sl, x, row_y, col_w, row_h, fill=bg)
            add_rect(sl, x, row_y, Inches(0.4), row_h, fill=C_ACCENT)
            add_text_box(sl, num, x, row_y, Inches(0.4), row_h,
                         font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_text_box(sl, label, x + Inches(0.48), row_y + Pt(4),
                         col_w - Inches(0.9), row_h - Pt(8),
                         font_size=13, color=C_DARK)
            add_text_box(sl, pg, x + col_w - Inches(0.75), row_y + Pt(4),
                         Inches(0.7), row_h - Pt(8),
                         font_size=12, color=C_MUTED, align=PP_ALIGN.RIGHT)

    draw_toc_col(items_col1, col1_x)
    draw_toc_col(items_col2, col2_x)
    return sl


# --- スライド 3: システム概要 ---
def slide_overview():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "システム概要・目的", "System Overview")
    draw_footer(sl, 3)

    Y = content_y()
    # 課題・目的カード
    card_w = Inches(5.8)
    card_h = Inches(2.1)

    # 課題カード
    add_rect(sl, Inches(0.4), Y, card_w, card_h, fill=C_WHITE, line=RGBColor(0xE0,0xE8,0xF0), line_w=Pt(1))
    add_rect(sl, Inches(0.4), Y, card_w, Inches(0.38), fill=RGBColor(0xC0,0x39,0x2B))
    add_text_box(sl, "▶ 解決した課題", Inches(0.5), Y + Pt(4), card_w - Inches(0.1), Inches(0.3),
                 font_size=13, bold=True, color=C_WHITE)
    txb = sl.shapes.add_textbox(Inches(0.5), Y + Inches(0.46), card_w - Inches(0.2), card_h - Inches(0.52))
    tf = txb.text_frame; tf.word_wrap = True
    for item in [
        "チャットでの連絡が直前になりがちで管理できなかった",
        "欠席の公平な評価ができず、選考に不満が出ていた",
        "練習参加状況を数値で把握する手段がなかった",
    ]:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = "・" + item
        r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.color.rgb = C_DARK
        p.space_before = Pt(3)

    # 目的カード
    add_rect(sl, Inches(6.8), Y, card_w, card_h, fill=C_WHITE, line=RGBColor(0xE0,0xE8,0xF0), line_w=Pt(1))
    add_rect(sl, Inches(6.8), Y, card_w, Inches(0.38), fill=C_ACCENT)
    add_text_box(sl, "▶ システムの目的", Inches(6.9), Y + Pt(4), card_w - Inches(0.1), Inches(0.3),
                 font_size=13, bold=True, color=C_WHITE)
    txb2 = sl.shapes.add_textbox(Inches(6.9), Y + Inches(0.46), card_w - Inches(0.2), card_h - Inches(0.52))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    for item in [
        "出欠連絡を火曜23:59までの事前申告に統一・徹底",
        "ポイントと技術ランクによる公平な選考スコアを算出",
        "管理者の手間を削減しながら規律を自動で強化する",
    ]:
        p = tf2.add_paragraph()
        r = p.add_run(); r.text = "・" + item
        r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.color.rgb = C_DARK
        p.space_before = Pt(3)

    # 主要機能サマリー
    add_rect(sl, Inches(0.4), Y + card_h + Inches(0.18), SLIDE_W - Inches(0.8), Inches(0.36), fill=C_NAV)
    add_text_box(sl, "主な機能", Inches(0.5), Y + card_h + Inches(0.18) + Pt(4),
                 Inches(2), Inches(0.28), font_size=13, bold=True, color=C_WHITE)

    funcs = [
        ("📅", "出欠登録"),
        ("📊", "ランキング"),
        ("🔔", "LINE通知"),
        ("📆", "Googleカレンダー連携"),
        ("👥", "メンバー管理"),
        ("⚠️", "警告フラグ"),
    ]
    feat_y = Y + card_h + Inches(0.62)
    feat_w = Inches(2.0)
    for i, (ic, lb) in enumerate(funcs):
        fx = Inches(0.4) + feat_w * i + Inches(0.05) * i
        add_rect(sl, fx, feat_y, feat_w, Inches(0.75), fill=C_WHITE, line=C_LINE, line_w=Pt(1))
        add_text_box(sl, ic, fx, feat_y + Pt(2), feat_w, Inches(0.35),
                     font_size=20, align=PP_ALIGN.CENTER)
        add_text_box(sl, lb, fx, feat_y + Inches(0.38), feat_w, Inches(0.35),
                     font_size=11, color=C_DARK, align=PP_ALIGN.CENTER)

    # スクショ
    add_screenshot_placeholder(sl, Inches(0.4), Y + card_h + Inches(1.55),
                                SLIDE_W - Inches(0.8), Inches(1.6))
    return sl


# --- スライド 4: ログイン ---
def slide_login():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "ログイン方法", "Login")
    draw_footer(sl, 4)

    Y = content_y()
    # LINEログイン手順
    steps = [
        ("1", "QRコードまたはURLでアプリにアクセス",
              "スマートフォンのLINEアプリが必要です"),
        ("2", "「LINEでログイン」ボタンをタップ",
              "LINE公式の認証画面が表示されます"),
        ("3", "LINEアカウントで認証を許可",
              "初回のみ情報提供への同意が必要です"),
        ("4", "氏名・学籍番号などを入力して申請",
              "幹部の承認後にシステムが使用できます"),
    ]
    step_w = Inches(2.9)
    step_h = Inches(2.0)
    for i, (num, title, sub) in enumerate(steps):
        sx = Inches(0.4) + (step_w + Inches(0.15)) * i
        add_rect(sl, sx, Y, step_w, step_h, fill=C_WHITE, line=C_LINE, line_w=Pt(1))
        # 丸数字
        add_rect(sl, sx + step_w/2 - Inches(0.27), Y - Inches(0.27),
                 Inches(0.54), Inches(0.54), fill=C_ACCENT)
        add_text_box(sl, num, sx + step_w/2 - Inches(0.27), Y - Inches(0.27),
                     Inches(0.54), Inches(0.54), font_size=16, bold=True,
                     color=C_WHITE, align=PP_ALIGN.CENTER)
        # 矢印（最後以外）
        if i < len(steps) - 1:
            add_text_box(sl, "▶", sx + step_w + Inches(0.02), Y + step_h/2 - Inches(0.18),
                         Inches(0.14), Inches(0.36), font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(sl, title, sx + Inches(0.12), Y + Inches(0.2), step_w - Inches(0.24), Inches(0.8),
                     font_size=13, bold=True, color=C_DARK, wrap=True)
        add_text_box(sl, sub, sx + Inches(0.12), Y + Inches(1.05), step_w - Inches(0.24), Inches(0.75),
                     font_size=11, color=C_MUTED, wrap=True)

    # 注意事項
    note_y = Y + step_h + Inches(0.2)
    add_rect(sl, Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.76),
             fill=RGBColor(0xFF, 0xF5, 0xE0), line=RGBColor(0xF3,0x9C,0x12), line_w=Pt(1.5))
    add_text_box(sl, "⚠ 注意事項",
                 Inches(0.55), note_y + Pt(4), Inches(1.5), Inches(0.3),
                 font_size=12, bold=True, color=RGBColor(0xD4,0x7F,0x00))
    add_text_box(sl, "承認前は「承認待ち」画面が表示され、各機能は利用できません。承認は幹部が行います。承認後にLINEで通知が届きます。",
                 Inches(0.55), note_y + Inches(0.32), SLIDE_W - Inches(1.2), Inches(0.38),
                 font_size=11, color=C_DARK, wrap=True)

    # スクショ
    add_screenshot_placeholder(sl, Inches(0.4), note_y + Inches(0.9), SLIDE_W - Inches(0.8), Inches(2.3))
    return sl


# --- スライド 5: ロール体系 ---
def slide_roles():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "ロール体系", "User Roles")
    draw_footer(sl, 5)

    Y = content_y()
    roles_data = [
        ("部　員", "member",     C_MEMBER,
         ["出欠連絡の登録・変更", "自分の成績・ポイント確認", "ランキング閲覧", "意見箱への投稿"]),
        ("マネージャー", "manager", C_MANAGER,
         ["部員と同じ全機能", "未提出者へのLINE通知送信", "出欠状況の一覧確認", "実績の確定・取消"]),
        ("幹　部", "admin",      C_ADMIN,
         ["マネージャーと同じ全機能", "メンバー承認・削除・ロール変更", "技術ランク編集", "警告フラグ管理", "意見箱の閲覧"]),
        ("顧　問", "coach",      C_COACH,
         ["出欠状況の閲覧のみ", "個人スコアは非表示", "練習一覧の参照"]),
    ]

    card_w = Inches(2.95)
    card_h = Inches(4.0)
    for i, (label, role_id, color, perms) in enumerate(roles_data):
        cx = Inches(0.35) + (card_w + Inches(0.12)) * i
        # カード枠
        add_rect(sl, cx, Y, card_w, card_h, fill=C_WHITE, line=color, line_w=Pt(2))
        # ヘッダー
        add_rect(sl, cx, Y, card_w, Inches(0.65), fill=color)
        add_text_box(sl, label, cx, Y + Pt(6), card_w, Inches(0.5),
                     font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 権限リスト
        txb = sl.shapes.add_textbox(cx + Inches(0.15), Y + Inches(0.75),
                                     card_w - Inches(0.3), card_h - Inches(0.85))
        tf = txb.text_frame; tf.word_wrap = True
        for p_text in perms:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = "✓  " + p_text
            r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.color.rgb = C_DARK
            p.space_before = Pt(5)

    # 凡例
    add_text_box(sl, "※ ロールは幹部のみが変更できます。承認直後のロールは「部員」です。",
                 Inches(0.4), Y + card_h + Inches(0.15), SLIDE_W - Inches(0.8), Inches(0.35),
                 font_size=11, color=C_MUTED)
    return sl


# --- スライド 6: ダッシュボード ---
def slide_dashboard():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "ダッシュボード", "Dashboard", role_color=C_MEMBER)
    draw_footer(sl, 6)
    add_role_badge(sl, "全ロール共通", C_ACCENT, SLIDE_W - Inches(2.1), Inches(0.36), Inches(1.9), Inches(0.38))

    Y = content_y()
    left_w = Inches(6.2)
    right_w = SLIDE_W - left_w - Inches(0.6)

    # 左：機能説明
    add_rect(sl, Inches(0.4), Y, left_w, SLIDE_H - Y - Inches(0.5), fill=C_WHITE, line=C_LINE, line_w=Pt(1))
    add_text_box(sl, "表示される情報", Inches(0.55), Y + Inches(0.1), left_w - Inches(0.2), Inches(0.38),
                 font_size=14, bold=True, color=C_NAV)
    add_rect(sl, Inches(0.55), Y + Inches(0.52), left_w - Inches(0.3), Pt(1.5), fill=C_LINE)

    items = [
        ("今日の練習情報", [
            "参加予定者一覧（出席 / 遅刻 / 欠席を色分け）",
            "欠席連絡者の理由確認",
            "自分の今日の出欠ステータス",
        ]),
        ("自分の出席状況カード", [
            "出席率（パーセンテージ＋グラフ）",
            "直近10回の活動履歴ドット",
            "無連絡欠席の警告件数",
        ]),
        ("ランキング", [
            "全体・学年別タブで切り替え可能",
            "選考スコア・出席率・ポイントを表示",
        ]),
    ]
    ty = Y + Inches(0.65)
    for group, subitems in items:
        add_text_box(sl, "▶ " + group, Inches(0.6), ty, left_w - Inches(0.3), Inches(0.35),
                     font_size=13, bold=True, color=C_ACCENT)
        ty += Inches(0.38)
        for sub in subitems:
            add_text_box(sl, "　・" + sub, Inches(0.65), ty, left_w - Inches(0.35), Inches(0.3),
                         font_size=11.5, color=C_DARK)
            ty += Inches(0.3)
        ty += Inches(0.1)

    # 右：スクショ
    rx = Inches(0.4) + left_w + Inches(0.2)
    add_screenshot_placeholder(sl, rx, Y, right_w, SLIDE_H - Y - Inches(0.5))
    return sl


# --- スライド 7: カレンダー・出欠登録 ---
def slide_calendar():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "カレンダー・出欠登録", "Calendar & Attendance", role_color=C_MEMBER)
    draw_footer(sl, 7)
    add_role_badge(sl, "全ロール共通", C_ACCENT, SLIDE_W - Inches(2.1), Inches(0.36), Inches(1.9), Inches(0.38))

    Y = content_y()
    # 上段：手順
    add_text_box(sl, "出欠登録の手順", Inches(0.4), Y, Inches(6), Inches(0.35),
                 font_size=14, bold=True, color=C_NAV)

    steps = [
        ("① カレンダーを開く", "メニューから「カレンダー」を選択"),
        ("② 練習日をタップ", "当該日程の詳細パネルが開く"),
        ("③ 出欠状況を選択", "出席 / 遅刻 / 欠席 から選ぶ"),
        ("④ 欠席の場合は理由を入力", "理由・詳細・遅刻時間帯を入力"),
        ("⑤ 送信ボタンで確定", "送信完了でカレンダーに反映"),
    ]
    step_w = Inches(2.4)
    for i, (s, d) in enumerate(steps):
        sx = Inches(0.4) + step_w * (i % 5) + Inches(0.08) * (i % 5)
        sy = Y + Inches(0.42)
        add_rect(sl, sx, sy, step_w, Inches(0.98), fill=C_WHITE, line=C_LINE, line_w=Pt(1))
        add_rect(sl, sx, sy, step_w, Inches(0.35), fill=C_ACCENT)
        add_text_box(sl, s, sx + Inches(0.05), sy + Pt(4), step_w - Inches(0.1), Inches(0.28),
                     font_size=11, bold=True, color=C_WHITE)
        add_text_box(sl, d, sx + Inches(0.08), sy + Inches(0.38), step_w - Inches(0.15), Inches(0.55),
                     font_size=10.5, color=C_DARK, wrap=True)

    # 欠席理由リスト
    reasons_y = Y + Inches(1.6)
    add_text_box(sl, "選択できる欠席理由", Inches(0.4), reasons_y, Inches(4), Inches(0.35),
                 font_size=13, bold=True, color=C_NAV)
    reasons = [
        ("体調不良", "翌日の参加がロック"),
        ("授業・補講", "学業優先"),
        ("その他", "詳細を自由記述"),
    ]
    for i, (r, note) in enumerate(reasons):
        rx_ = Inches(0.4) + Inches(4.15) * i
        add_rect(sl, rx_, reasons_y + Inches(0.4), Inches(3.9), Inches(0.65),
                 fill=C_WHITE, line=C_LINE, line_w=Pt(1))
        add_text_box(sl, r, rx_ + Inches(0.1), reasons_y + Inches(0.44),
                     Inches(2.4), Inches(0.3), font_size=12, bold=True, color=C_DARK)
        add_text_box(sl, note, rx_ + Inches(0.1), reasons_y + Inches(0.76),
                     Inches(3.6), Inches(0.28), font_size=10, color=C_MUTED)

    # 遅刻時の限数選択
    add_text_box(sl, "遅刻の場合は到着できる授業の限数を選択（8・9・10限）",
                 Inches(0.4), reasons_y + Inches(1.2), SLIDE_W - Inches(0.8), Inches(0.3),
                 font_size=11.5, color=RGBColor(0x1A,0x7A,0x3C))

    # スクショ
    add_screenshot_placeholder(sl, Inches(0.4), reasons_y + Inches(1.6),
                                SLIDE_W - Inches(0.8), Inches(1.9))
    return sl


# --- スライド 8: 出欠連絡のルール ---
def slide_rules():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "出欠連絡のルール（締め切り）", "Attendance Rules", role_color=C_MEMBER)
    draw_footer(sl, 8)

    Y = content_y()
    add_text_box(sl, "出欠連絡は毎週 火曜日 23:59 が締め切りです",
                 Inches(0.4), Y, SLIDE_W - Inches(0.8), Inches(0.42),
                 font_size=16, bold=True, color=C_NAV)
    add_text_box(sl, "水〜金曜の練習に対して、前週土曜〜当週火曜の間に連絡を済ませてください。",
                 Inches(0.4), Y + Inches(0.45), SLIDE_W - Inches(0.8), Inches(0.35),
                 font_size=13, color=C_DARK)

    # 週間カレンダー図
    days = ["土", "日", "月", "火", "水", "木", "金"]
    day_w = Inches(1.7)
    day_h = Inches(1.3)
    day_y = Y + Inches(0.9)
    for i, d in enumerate(days):
        dx = Inches(0.4) + day_w * i + Inches(0.04) * i
        is_deadline = (d == "火")
        is_ok = (d in ["土", "日", "月", "火"])
        is_practice = (d in ["水", "木", "金"])
        bg = RGBColor(0xE8,0xF5,0xE9) if is_ok else (RGBColor(0xE3,0xF2,0xFD) if is_practice else C_WHITE)
        border = C_MEMBER if is_ok else (C_ACCENT if is_practice else C_LINE)
        border_w = Pt(2.5) if (is_deadline or is_practice) else Pt(1)
        add_rect(sl, dx, day_y, day_w, day_h, fill=bg, line=border, line_w=border_w)
        add_text_box(sl, d, dx, day_y + Pt(6), day_w, Inches(0.4),
                     font_size=18, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        if is_deadline:
            add_text_box(sl, "⏰ 締め切り\n23:59", dx, day_y + Inches(0.5), day_w, Inches(0.7),
                         font_size=10, bold=True, color=RGBColor(0xC0,0x39,0x2B), align=PP_ALIGN.CENTER)
        elif is_ok:
            add_text_box(sl, "✓ 連絡OK", dx, day_y + Inches(0.5), day_w, Inches(0.45),
                         font_size=10, color=C_MEMBER, align=PP_ALIGN.CENTER)
        elif is_practice:
            add_text_box(sl, "練習日", dx, day_y + Inches(0.5), day_w, Inches(0.45),
                         font_size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)
            add_text_box(sl, "連絡受付なし", dx, day_y + Inches(0.8), day_w, Inches(0.4),
                         font_size=9, color=C_MUTED, align=PP_ALIGN.CENTER)

    # 注意事項ボックス
    warn_y = day_y + day_h + Inches(0.2)
    add_rect(sl, Inches(0.4), warn_y, SLIDE_W - Inches(0.8), Inches(1.1),
             fill=RGBColor(0xFF,0xF3,0xE0), line=RGBColor(0xF5,0x9A,0x00), line_w=Pt(1.5))
    add_text_box(sl, "⚠  重要な制約", Inches(0.6), warn_y + Pt(5), Inches(2), Inches(0.3),
                 font_size=12, bold=True, color=RGBColor(0xD4,0x7F,0x00))
    warn_items = [
        "水〜金曜の練習当日は連絡を受け付けません（当日欠席でも-50点）",
        "当日欠席とは：練習開始1時間以内の連絡（-50点）",
        "無連絡欠席は最も重い -100点 のペナルティです",
    ]
    for j, wi in enumerate(warn_items):
        add_text_box(sl, "・" + wi, Inches(0.6), warn_y + Inches(0.36) + Inches(0.24)*j,
                     SLIDE_W - Inches(1.2), Inches(0.25),
                     font_size=11, color=C_DARK)

    # スクショ
    add_screenshot_placeholder(sl, Inches(0.4), warn_y + Inches(1.2), SLIDE_W - Inches(0.8), Inches(1.6))
    return sl


# --- スライド 9: ポイントシステム ---
def slide_points():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "ポイントシステム", "Point System", role_color=C_MEMBER)
    draw_footer(sl, 9)
    add_role_badge(sl, "全ロール共通", C_ACCENT, SLIDE_W - Inches(2.1), Inches(0.36), Inches(1.9), Inches(0.38))

    Y = content_y()
    add_text_box(sl, "全員が 1,000 点からスタートし、出欠状況に応じてポイントが変動します。",
                 Inches(0.4), Y, SLIDE_W - Inches(0.8), Inches(0.35), font_size=13, color=C_DARK)

    # ポイントテーブル
    table_y = Y + Inches(0.45)
    headers = ["ステータス", "条件", "ポイント変動"]
    col_ws = [Inches(3.5), Inches(6.0), Inches(2.5)]
    row_h = Inches(0.52)

    def table_row(row_y, cells, bg, text_colors=None):
        x = Inches(0.4)
        for ci, (cell, cw) in enumerate(zip(cells, col_ws)):
            add_rect(sl, x, row_y, cw, row_h, fill=bg, line=C_LINE, line_w=Pt(0.5))
            tc = (text_colors[ci] if text_colors else C_DARK)
            add_text_box(sl, cell, x + Inches(0.1), row_y + Pt(5), cw - Inches(0.15), row_h - Pt(8),
                         font_size=12, bold=(bg == C_NAV), color=tc, wrap=True)
            x += cw

    table_row(table_y, headers, C_NAV, [C_WHITE, C_WHITE, C_WHITE])
    rows = [
        ("出席",           "当日参加",                          "±0"),
        ("遅刻",           "予定より遅れて参加",                  "−10点"),
        ("通常欠席",        "1時間以上前に連絡した欠席",            "−20点"),
        ("当日欠席",        "練習開始 1時間以内に連絡した欠席",      "−50点"),
        ("無連絡欠席",      "連絡なしで欠席（実績確定時に確定）",     "−100点"),
    ]
    bgs = [C_WHITE, RGBColor(0xF2,0xF7,0xFF)] * 5
    point_colors = [
        C_DARK,
        RGBColor(0xD4,0x7F,0x00),
        RGBColor(0xD4,0x7F,0x00),
        RGBColor(0xC0,0x39,0x2B),
        RGBColor(0xC0,0x39,0x2B),
    ]
    for i, (status, cond, pts) in enumerate(rows):
        tc = [C_DARK, C_DARK, point_colors[i]]
        table_row(table_y + row_h * (i+1), [status, cond, pts], bgs[i], tc)

    # スクショ
    note_y = table_y + row_h * 6 + Inches(0.15)
    add_text_box(sl, "※ ポイントはダッシュボードのカードから確認できます。過去の変動履歴も参照可能です。",
                 Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.3), font_size=11, color=C_MUTED)
    add_screenshot_placeholder(sl, Inches(0.4), note_y + Inches(0.35), SLIDE_W - Inches(0.8), Inches(1.85))
    return sl


# --- スライド 10: ランキング ---
def slide_ranking():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "ランキング・選考スコア", "Ranking", role_color=C_MEMBER)
    draw_footer(sl, 10)
    add_role_badge(sl, "全ロール共通", C_ACCENT, SLIDE_W - Inches(2.1), Inches(0.36), Inches(1.9), Inches(0.38))

    Y = content_y()
    left_x = Inches(0.4)
    left_w = Inches(6.2)
    right_x = left_x + left_w + Inches(0.2)
    right_w = SLIDE_W - right_x - Inches(0.2)

    # 左：説明
    add_text_box(sl, "選考スコアの算出式", left_x, Y, left_w, Inches(0.35),
                 font_size=14, bold=True, color=C_NAV)
    # 数式ボックス
    add_rect(sl, left_x, Y + Inches(0.4), left_w, Inches(0.75), fill=C_NAV)
    add_text_box(sl, "選考スコア  ＝  ( 技術ランク ÷ 6 )  ×  出席率  ×  100",
                 left_x + Inches(0.15), Y + Inches(0.48), left_w - Inches(0.3), Inches(0.58),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(sl, "出席率 ＝ ( 出席数 ＋ 遅刻数 × 0.5 ) ÷ 総セッション数 × 100",
                 left_x, Y + Inches(1.25), left_w, Inches(0.3),
                 font_size=11.5, color=C_DARK)

    # ランク区分テーブル
    add_text_box(sl, "選考ランク区分", left_x, Y + Inches(1.65), left_w, Inches(0.32),
                 font_size=13, bold=True, color=C_NAV)
    ranks = [("S", "85以上",  RGBColor(0xFF,0xD7,0x00)),
             ("A", "70以上",  RGBColor(0xC0,0x39,0x2B)),
             ("B", "55以上",  RGBColor(0xD4,0x7F,0x00)),
             ("C", "40以上",  RGBColor(0x27,0xAE,0x60)),
             ("D", "25以上",  C_ACCENT),
             ("E", "25未満",  C_MUTED)]
    rk_w = Inches(0.95)
    for i, (r, cond, rc) in enumerate(ranks):
        rx_ = left_x + rk_w * i + Inches(0.05) * i
        ry_ = Y + Inches(2.0)
        add_rect(sl, rx_, ry_, rk_w, Inches(0.55), fill=rc)
        add_text_box(sl, r, rx_, ry_, rk_w, Inches(0.35), font_size=18, bold=True,
                     color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, cond, rx_, ry_ + Inches(0.33), rk_w, Inches(0.24), font_size=9,
                     color=C_WHITE, align=PP_ALIGN.CENTER)

    # 技術ランク説明
    add_text_box(sl, "技術ランクについて", left_x, Y + Inches(2.68), left_w, Inches(0.32),
                 font_size=13, bold=True, color=C_NAV)
    sk_ranks = [("E","1"),("D","2"),("C","3"),("B","4"),("A","5"),("S","6")]
    for i, (r, v) in enumerate(sk_ranks):
        srx = left_x + Inches(0.95) * i + Inches(0.05)*i
        add_rect(sl, srx, Y + Inches(3.0), Inches(0.9), Inches(0.55),
                 fill=C_WHITE, line=C_LINE, line_w=Pt(1))
        add_text_box(sl, r, srx, Y + Inches(3.0), Inches(0.9), Inches(0.3),
                     font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(sl, f"Lv.{v}", srx, Y + Inches(3.3), Inches(0.9), Inches(0.25),
                     font_size=10, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_text_box(sl, "初期ランクは C（Lv.3）。幹部のみ変更可。",
                 left_x, Y + Inches(3.65), left_w, Inches(0.3), font_size=11, color=C_MUTED)

    # 右：スクショ
    add_screenshot_placeholder(sl, right_x, Y, right_w, SLIDE_H - Y - Inches(0.5))
    return sl


# --- スライド 11: マネージャー機能 ---
def slide_manager():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "マネージャー専用機能", "Manager Features", role_color=C_MANAGER)
    draw_footer(sl, 11)
    add_role_badge(sl, "マネージャー以上", C_MANAGER, SLIDE_W - Inches(2.5), Inches(0.36), Inches(2.3), Inches(0.38))

    Y = content_y()
    add_text_box(sl, "マネージャーは部員と同じ機能に加え、以下の管理機能が使用できます。",
                 Inches(0.4), Y, SLIDE_W - Inches(0.8), Inches(0.35), font_size=13, color=C_DARK)

    funcs = [
        ("未提出者へのLINE通知", [
            "指定した練習日の出欠を未提出の部員に一括通知",
            "通知はLINEメッセージで個別送信",
            "送信件数・対象者を画面で確認してから実行",
        ]),
        ("出欠状況の一覧管理", [
            "全部員の出欠状況を練習日ごとに確認",
            "未提出者・欠席者の一覧表示",
        ]),
        ("実績確定", [
            "練習終了後に管理者が出欠を「確定」する操作",
            "確定後にポイントが反映・スコア計算に使用",
            "無連絡欠席はここで管理者が手動設定",
        ]),
    ]

    fy = Y + Inches(0.45)
    for i, (title, items) in enumerate(funcs):
        fy_start = fy + Inches(1.55) * i
        add_rect(sl, Inches(0.4), fy_start, SLIDE_W * 0.55, Inches(1.45),
                 fill=C_WHITE, line=C_MANAGER, line_w=Pt(1.5))
        add_rect(sl, Inches(0.4), fy_start, SLIDE_W * 0.55, Inches(0.38), fill=C_MANAGER)
        add_text_box(sl, title, Inches(0.55), fy_start + Pt(4),
                     SLIDE_W * 0.55 - Inches(0.2), Inches(0.3),
                     font_size=13, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            add_text_box(sl, "　・" + item,
                         Inches(0.55), fy_start + Inches(0.44) + Inches(0.3) * j,
                         SLIDE_W * 0.55 - Inches(0.2), Inches(0.3),
                         font_size=11.5, color=C_DARK)

    # 右：スクショ
    rx = Inches(0.4) + SLIDE_W * 0.55 + Inches(0.2)
    rw = SLIDE_W - rx - Inches(0.2)
    add_screenshot_placeholder(sl, rx, Y + Inches(0.45), rw, SLIDE_H - Y - Inches(1.0))
    return sl


# --- スライド 12: 幹部：メンバー管理 ---
def slide_admin_members():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "幹部：メンバー管理", "Admin - Member Management", role_color=C_ADMIN)
    draw_footer(sl, 12)
    add_role_badge(sl, "幹部のみ", C_ADMIN, SLIDE_W - Inches(2.0), Inches(0.36), Inches(1.8), Inches(0.38))

    Y = content_y()
    left_w = Inches(6.2)
    right_x = Inches(0.4) + left_w + Inches(0.2)
    right_w = SLIDE_W - right_x - Inches(0.2)

    tasks = [
        ("入部申請の承認・却下", [
            "新規ユーザーの承認待ちリストを確認",
            "承認するとロール「部員」で登録完了",
            "却下した場合は申請が削除される",
        ]),
        ("プロフィール編集", [
            "表示名・学年・ロールを変更可能",
            "技術ランク（E〜S）を幹部が設定",
            "ロール変更：部員 → マネージャー → 幹部",
        ]),
        ("メンバー削除", [
            "アカウントと出欠データを完全削除",
            "確認ダイアログで誤操作を防止",
        ]),
    ]

    ty = Y
    for title, items in tasks:
        add_rect(sl, Inches(0.4), ty, left_w, Inches(0.34), fill=C_ADMIN)
        add_text_box(sl, title, Inches(0.55), ty + Pt(3), left_w - Inches(0.2), Inches(0.28),
                     font_size=13, bold=True, color=C_WHITE)
        ty += Inches(0.34)
        bh = Inches(0.33) * len(items)
        add_rect(sl, Inches(0.4), ty, left_w, bh, fill=C_WHITE, line=C_LINE, line_w=Pt(0.5))
        for j, item in enumerate(items):
            add_text_box(sl, "　・" + item, Inches(0.55), ty + Inches(0.06) + Inches(0.33)*j,
                         left_w - Inches(0.2), Inches(0.3), font_size=11.5, color=C_DARK)
        ty += bh + Inches(0.18)

    add_screenshot_placeholder(sl, right_x, Y, right_w, SLIDE_H - Y - Inches(0.5))
    return sl


# --- スライド 13: 幹部：実績確定 ---
def slide_admin_confirm():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "幹部：実績確定", "Admin - Results Confirmation", role_color=C_ADMIN)
    draw_footer(sl, 13)
    add_role_badge(sl, "幹部・マネージャー", C_MANAGER, SLIDE_W - Inches(2.8), Inches(0.36), Inches(2.6), Inches(0.38))

    Y = content_y()
    add_text_box(sl, "実績確定とは？",
                 Inches(0.4), Y, Inches(5), Inches(0.35), font_size=14, bold=True, color=C_NAV)
    add_text_box(sl,
                 "部員の自己申告（事前登録）に対して、管理者が練習後に最終ステータスを確定する操作です。\n"
                 "確定後のデータがポイント計算・選考スコアに使用されます。",
                 Inches(0.4), Y + Inches(0.38), SLIDE_W - Inches(0.8), Inches(0.58),
                 font_size=12, color=C_DARK, wrap=True)

    # 2段階フロー
    flow_y = Y + Inches(1.1)
    add_text_box(sl, "2段階の確定フロー", Inches(0.4), flow_y, Inches(6), Inches(0.32),
                 font_size=13, bold=True, color=C_NAV)
    boxes = [
        ("① 部員が事前登録", "出席 / 遅刻 / 欠席 を火曜23:59までに入力", C_MEMBER),
        ("→", "", C_WHITE),
        ("② 管理者が実績確定", "練習後に確認・ステータスを確定\n無連絡欠席は管理者が手動設定", C_ADMIN),
        ("→", "", C_WHITE),
        ("③ スコアへ反映", "確定済みデータでポイント・\n選考スコアを自動計算", C_ACCENT),
    ]
    bx = Inches(0.4)
    for i, (t, d, c) in enumerate(boxes):
        bw = Inches(3.8) if t != "→" else Inches(0.4)
        by = flow_y + Inches(0.38)
        if t == "→":
            add_text_box(sl, "▶", bx, by + Inches(0.3), bw, Inches(0.4),
                         font_size=16, color=C_ACCENT, align=PP_ALIGN.CENTER)
        else:
            add_rect(sl, bx, by, bw, Inches(1.05), fill=C_WHITE, line=c, line_w=Pt(2))
            add_rect(sl, bx, by, bw, Inches(0.35), fill=c)
            add_text_box(sl, t, bx + Inches(0.1), by + Pt(4), bw - Inches(0.15), Inches(0.28),
                         font_size=12, bold=True, color=C_WHITE)
            add_text_box(sl, d, bx + Inches(0.1), by + Inches(0.4), bw - Inches(0.15), Inches(0.62),
                         font_size=11, color=C_DARK, wrap=True)
        bx += bw + Inches(0.05)

    # 注意
    note_y = flow_y + Inches(1.6)
    add_rect(sl, Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.6),
             fill=RGBColor(0xF3,0xE8,0xFF), line=C_ADMIN, line_w=Pt(1.5))
    add_text_box(sl, "⚠ 注意",
                 Inches(0.6), note_y + Pt(5), Inches(1), Inches(0.28),
                 font_size=12, bold=True, color=C_ADMIN)
    add_text_box(sl, "確定後の変更も可能ですが、ポイント再計算は管理者の手動操作が必要です。慎重に確定してください。",
                 Inches(0.6), note_y + Inches(0.3), SLIDE_W - Inches(1.2), Inches(0.25),
                 font_size=11, color=C_DARK)

    add_screenshot_placeholder(sl, Inches(0.4), note_y + Inches(0.75), SLIDE_W - Inches(0.8), Inches(2.3))
    return sl


# --- スライド 14: 幹部：警告フラグ ---
def slide_warning():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "幹部：警告フラグ管理", "Admin - Warning Flags", role_color=C_ADMIN)
    draw_footer(sl, 14)
    add_role_badge(sl, "幹部のみ", C_ADMIN, SLIDE_W - Inches(2.0), Inches(0.36), Inches(1.8), Inches(0.38))

    Y = content_y()
    add_text_box(sl, "部員に対して3種類の警告フラグを設定できます。フラグは段階的に重大度が上がります。",
                 Inches(0.4), Y, SLIDE_W - Inches(0.8), Inches(0.35), font_size=12.5, color=C_DARK)

    flags = [
        ("部費滞納", "dues_overdue",     "部費の未払いが続く場合"),
        ("無連絡欠席継続", "absent_no_report", "無連絡欠席が繰り返される場合"),
        ("品位違反", "conduct",          "チームのルール・マナーを違反した場合"),
    ]
    severities = [
        ("warning",       "警告",     RGBColor(0xF3,0x9C,0x12)),
        ("final_warning", "最終警告", RGBColor(0xE7,0x4C,0x3C)),
        ("expelled",      "除籍",     RGBColor(0x6C,0x35,0x7B)),
    ]

    # フラグ種類
    add_text_box(sl, "フラグの種類", Inches(0.4), Y + Inches(0.45), Inches(6), Inches(0.3),
                 font_size=13, bold=True, color=C_NAV)
    for i, (name, code, desc) in enumerate(flags):
        fy = Y + Inches(0.82) + Inches(0.62) * i
        add_rect(sl, Inches(0.4), fy, Inches(6.0), Inches(0.56),
                 fill=C_WHITE, line=RGBColor(0xE7,0x4C,0x3C), line_w=Pt(1))
        add_rect(sl, Inches(0.4), fy, Inches(0.06), Inches(0.56), fill=RGBColor(0xE7,0x4C,0x3C))
        add_text_box(sl, name, Inches(0.55), fy + Pt(3), Inches(2.0), Inches(0.28),
                     font_size=13, bold=True, color=C_DARK)
        add_text_box(sl, desc, Inches(0.55), fy + Inches(0.3), Inches(5.7), Inches(0.22),
                     font_size=10.5, color=C_MUTED)

    # 重大度レベル
    add_text_box(sl, "重大度レベル", Inches(6.7), Y + Inches(0.45), Inches(6), Inches(0.3),
                 font_size=13, bold=True, color=C_NAV)
    for i, (sev_id, label, color) in enumerate(severities):
        sy = Y + Inches(0.82) + Inches(0.65) * i
        add_rect(sl, Inches(6.7), sy, Inches(5.8), Inches(0.58),
                 fill=color)
        add_text_box(sl, label, Inches(6.85), sy + Pt(6), Inches(2.5), Inches(0.38),
                     font_size=15, bold=True, color=C_WHITE)
        note = ["部則の遵守を促す最初の警告", "改善なければ退部も検討する段階", "部からの除籍処分"][i]
        add_text_box(sl, note, Inches(9.0), sy + Pt(9), Inches(3.4), Inches(0.36),
                     font_size=11, color=C_WHITE)

    add_screenshot_placeholder(sl, Inches(0.4), Y + Inches(2.65), SLIDE_W - Inches(0.8), Inches(2.5))
    return sl


# --- スライド 15: 選考スコア詳細 ---
def slide_score_detail():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "選考スコアの仕組み", "Selection Score Detail")
    draw_footer(sl, 15)

    Y = content_y()
    # 算出例テーブル
    add_text_box(sl, "計算例", Inches(0.4), Y, Inches(5), Inches(0.35),
                 font_size=14, bold=True, color=C_NAV)

    examples = [
        ("部員A", "技術ランクB（Lv.4）", "出席率 85%", "4÷6 × 85 ≈ 56.7", "B"),
        ("部員B", "技術ランクS（Lv.6）", "出席率 70%", "6÷6 × 70 = 70.0",  "A"),
        ("部員C", "技術ランクC（Lv.3）", "出席率 90%", "3÷6 × 90 = 45.0",  "C"),
    ]

    # ヘッダー
    hdrs = ["部員", "技術ランク", "出席率", "計算", "選考ランク"]
    hcols = [Inches(1.5), Inches(2.5), Inches(1.5), Inches(4.5), Inches(1.8)]
    hx = Inches(0.4)
    for h, cw in zip(hdrs, hcols):
        add_rect(sl, hx, Y + Inches(0.4), cw, Inches(0.42), fill=C_NAV, line=C_LINE, line_w=Pt(0.5))
        add_text_box(sl, h, hx + Inches(0.05), Y + Inches(0.44), cw - Inches(0.1), Inches(0.35),
                     font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        hx += cw

    rank_colors = {"A": RGBColor(0xC0,0x39,0x2B), "B": RGBColor(0xD4,0x7F,0x00), "C": C_MEMBER}
    for ri, (name, rank, rate, calc, sel) in enumerate(examples):
        rx_ = Inches(0.4)
        bg = C_WHITE if ri % 2 == 0 else RGBColor(0xF2,0xF7,0xFF)
        ry = Y + Inches(0.82) + Inches(0.42) * ri
        for ci, (cell, cw) in enumerate(zip([name, rank, rate, calc, sel], hcols)):
            add_rect(sl, rx_, ry, cw, Inches(0.42), fill=bg, line=C_LINE, line_w=Pt(0.5))
            clr = rank_colors.get(sel, C_DARK) if ci == 4 else C_DARK
            bld = ci == 4
            add_text_box(sl, cell, rx_ + Inches(0.05), ry + Pt(5), cw - Inches(0.1), Inches(0.35),
                         font_size=12, bold=bld, color=clr, align=PP_ALIGN.CENTER)
            rx_ += cw

    # 注意点
    note_y = Y + Inches(2.15)
    add_text_box(sl, "スコアに関する注意事項", Inches(0.4), note_y, Inches(6), Inches(0.32),
                 font_size=13, bold=True, color=C_NAV)
    notes_list = [
        "スコアは確定済み（result_status）のセッションのみが対象",
        "出席回数が6回未満の部員はランキングからフィルタリング可能",
        "技術ランクは幹部のみが変更できます（初期値：C / Lv.3）",
        "顧問は個人スコアの数値が非表示で、ランキングを閲覧できます",
    ]
    for i, n in enumerate(notes_list):
        add_text_box(sl, "・" + n, Inches(0.5), note_y + Inches(0.38) + Inches(0.3)*i,
                     SLIDE_W - Inches(0.9), Inches(0.3), font_size=12, color=C_DARK)

    add_screenshot_placeholder(sl, Inches(0.4), note_y + Inches(1.65), SLIDE_W - Inches(0.8), Inches(2.1))
    return sl


# --- スライド 16: 体調不良ロック ---
def slide_sicklock():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "体調不良ロック機能", "Sick Lockout", role_color=C_MEMBER)
    draw_footer(sl, 16)

    Y = content_y()
    # 説明
    add_text_box(sl, "体調不良を理由に欠席連絡した場合、翌日の練習への参加登録がロックされます。",
                 Inches(0.4), Y, SLIDE_W - Inches(0.8), Inches(0.38), font_size=13, color=C_DARK)

    # フロー図
    flow_y = Y + Inches(0.5)
    flow_items = [
        ("体調不良で\n欠席を登録", C_MEMBER),
        ("翌日の\n練習がロック", RGBColor(0xE7,0x4C,0x3C)),
        ("ロック期間中は\n出欠登録不可", RGBColor(0xD4,0x7F,0x00)),
        ("翌日以降は\n通常通り利用可", C_ACCENT),
    ]
    fw = Inches(2.8)
    for i, (label, c) in enumerate(flow_items):
        fx = Inches(0.4) + (fw + Inches(0.3)) * i
        add_rect(sl, fx, flow_y, fw, Inches(1.1), fill=c)
        add_text_box(sl, label, fx, flow_y + Inches(0.2), fw, Inches(0.7),
                     font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < len(flow_items) - 1:
            add_text_box(sl, "▶", fx + fw + Inches(0.05), flow_y + Inches(0.35),
                         Inches(0.22), Inches(0.4), font_size=16, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # 目的ボックス
    purpose_y = flow_y + Inches(1.25)
    add_rect(sl, Inches(0.4), purpose_y, SLIDE_W - Inches(0.8), Inches(0.85),
             fill=RGBColor(0xE8,0xF4,0xFD), line=C_ACCENT, line_w=Pt(1.5))
    add_text_box(sl, "💡 この機能の目的", Inches(0.6), purpose_y + Pt(5), Inches(3), Inches(0.3),
                 font_size=13, bold=True, color=C_ACCENT)
    add_text_box(sl,
                 "「体調不良」を口実とした安易な欠席を抑止するためのルールです。\n"
                 "翌日に健康を回復しているなら練習に参加できるはずという観点から設計されています。",
                 Inches(0.6), purpose_y + Inches(0.35), SLIDE_W - Inches(1.2), Inches(0.46),
                 font_size=12, color=C_DARK, wrap=True)

    # ダッシュボード表示説明
    disp_y = purpose_y + Inches(1.0)
    add_text_box(sl, "ロック中のダッシュボード表示",
                 Inches(0.4), disp_y, Inches(5), Inches(0.3), font_size=13, bold=True, color=C_NAV)
    add_rect(sl, Inches(0.4), disp_y + Inches(0.35), SLIDE_W * 0.45, Inches(0.65),
             fill=RGBColor(0xFF,0xEB,0xEB), line=RGBColor(0xE7,0x4C,0x3C), line_w=Pt(1.5))
    add_text_box(sl, "⚠ 体調不良ロック中 ― 本日の練習参加登録は制限されています",
                 Inches(0.55), disp_y + Inches(0.42), SLIDE_W * 0.43, Inches(0.5),
                 font_size=11, bold=True, color=RGBColor(0xC0,0x30,0x20))

    add_screenshot_placeholder(sl, Inches(0.4), disp_y + Inches(1.1), SLIDE_W - Inches(0.8), Inches(2.35))
    return sl


# --- スライド 17: FAQ ---
def slide_faq():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl)
    draw_header(sl, "よくある質問（FAQ）", "Frequently Asked Questions")
    draw_footer(sl, 17)

    Y = content_y()
    faqs = [
        ("Q", "出欠を間違えて登録してしまいました。修正できますか？",
         "A", "火曜23:59の締め切り前であれば、カレンダーから再度同じ日をタップして変更できます。締め切り後の修正は幹部へご相談ください。"),
        ("Q", "ポイントが思っていたより減っていました。なぜですか？",
         "A", "ポイントは「実績確定」後に反映されます。確定前は表示が変わりません。また「当日欠席」は-50点です。"),
        ("Q", "ランキングに自分が表示されません。",
         "A", "出席回数が6回未満の場合、フィルタリングされて非表示になることがあります。フィルターをオフにすると表示できます。"),
        ("Q", "LINEの通知が届きません。",
         "A", "LINEのブロック設定や通知設定をご確認ください。また、LINEログインで登録したアカウントでのみ通知が届きます。"),
        ("Q", "技術ランクを変えたいのですが。",
         "A", "技術ランクは幹部のみが変更できます。変更を希望する場合は幹部にお申し出ください。"),
    ]

    for i, (ql, qt, al, at) in enumerate(faqs):
        fy = Y + Inches(1.14) * i
        # Q
        add_rect(sl, Inches(0.4), fy, Inches(0.48), Inches(0.42), fill=C_ACCENT)
        add_text_box(sl, ql, Inches(0.4), fy, Inches(0.48), Inches(0.42),
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(0.88), fy, SLIDE_W - Inches(1.28), Inches(0.42),
                 fill=RGBColor(0xE8,0xF1,0xFB), line=C_LINE, line_w=Pt(0.5))
        add_text_box(sl, qt, Inches(0.98), fy + Pt(4), SLIDE_W - Inches(1.38), Inches(0.35),
                     font_size=11.5, bold=True, color=C_NAV)
        # A
        ay = fy + Inches(0.44)
        add_rect(sl, Inches(0.4), ay, Inches(0.48), Inches(0.62), fill=C_MEMBER)
        add_text_box(sl, al, Inches(0.4), ay, Inches(0.48), Inches(0.62),
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(0.88), ay, SLIDE_W - Inches(1.28), Inches(0.62),
                 fill=C_WHITE, line=C_LINE, line_w=Pt(0.5))
        add_text_box(sl, at, Inches(0.98), ay + Pt(4), SLIDE_W - Inches(1.4), Inches(0.55),
                     font_size=11, color=C_DARK, wrap=True)

    return sl


# --- スライド 18: お問い合わせ ---
def slide_contact():
    sl = prs.slides.add_slide(BLANK)
    draw_base(sl, C_NAV)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=C_NAV)
    add_rect(sl, 0, Inches(3.5), SLIDE_W, Pt(3), fill=C_ACCENT)

    add_text_box(sl, "お問い合わせ・サポート",
                 Inches(1), Inches(0.8), SLIDE_W - Inches(2), Inches(0.7),
                 font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "システムの不具合・操作に関するご質問",
                 Inches(1), Inches(1.55), SLIDE_W - Inches(2), Inches(0.4),
                 font_size=15, color=RGBColor(0x80,0xB8,0xE8), align=PP_ALIGN.CENTER)

    contacts = [
        ("🏸", "幹部への直接相談", "LINEグループまたは直接ご連絡ください"),
        ("⚙️", "技術的な問題", "システム管理者（幹部）にスクリーンショット付きで報告"),
        ("📋", "ルール・判断の疑問", "練習前後に幹部に確認してください"),
    ]
    for i, (ic, title, desc) in enumerate(contacts):
        cx = Inches(0.8) + Inches(4.1) * i
        cy = Inches(4.0)
        add_rect(sl, cx, cy, Inches(3.8), Inches(1.8), fill=RGBColor(0x25,0x40,0x60))
        add_text_box(sl, ic, cx, cy + Inches(0.15), Inches(3.8), Inches(0.5),
                     font_size=24, align=PP_ALIGN.CENTER)
        add_text_box(sl, title, cx, cy + Inches(0.65), Inches(3.8), Inches(0.4),
                     font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, cx + Inches(0.1), cy + Inches(1.1), Inches(3.6), Inches(0.6),
                     font_size=10.5, color=RGBColor(0x80,0xB8,0xE8), align=PP_ALIGN.CENTER, wrap=True)

    add_text_box(sl, "千葉工業大学 バドミントン部  ―  BadAttend システム",
                 Inches(1), SLIDE_H - Inches(0.8), SLIDE_W - Inches(2), Inches(0.35),
                 font_size=11, color=RGBColor(0x50,0x70,0x90), align=PP_ALIGN.CENTER)
    return sl


# ===================== スライド生成 =====================
print("スライドを生成中...")
slide_title()
slide_toc()
slide_overview()
slide_login()
slide_roles()
slide_dashboard()
slide_calendar()
slide_rules()
slide_points()
slide_ranking()
slide_manager()
slide_admin_members()
slide_admin_confirm()
slide_warning()
slide_score_detail()
slide_sicklock()
slide_faq()
slide_contact()

OUT = r"c:\Users\kenta\project\BadAttend\BadAttend_マニュアル.pptx"
prs.save(OUT)
print(f"完了: {OUT}")
print(f"スライド数: {len(prs.slides)}")
